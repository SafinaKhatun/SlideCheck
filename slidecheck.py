import os
import re
from pptx import Presentation
from PIL import Image
import pytesseract

def extract_text_from_pptx(file_path):
    """Extract text from each slide in the PowerPoint presentation."""
    prs = Presentation(file_path)
    text_data = []
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        text_data.append(slide_text.strip())
    return text_data

def extract_text_from_image(image_path):
    """Extract text from an image using Tesseract OCR."""
    img = Image.open(image_path)
    return pytesseract.image_to_string(img)

def check_conflicting_numerical_data(text_data):
    """Identify conflicting numerical data across slides."""
    numerical_data = {}
    inconsistencies = []
    
    for i, text in enumerate(text_data):
         
        numbers = re.findall(r'(\d+\.?\d*%?)', text)
        for number in numbers:
             
            normalized_number = number.replace('%', '')
            if normalized_number in numerical_data:
                numerical_data[normalized_number].append(i + 1)
            else:
                numerical_data[normalized_number] = [i + 1]
    
    for number, slides in numerical_data.items():
        if len(slides) > 1:
            inconsistencies.append(f"Conflicting numerical data '{number}' found on slides: {slides}")
    
    return inconsistencies

def check_contradictory_claims(text_data):
    """Check for contradictory claims within the text."""
    contradictions = []
    claims = [
        ("market is highly competitive", "few competitors"),
        ("strong growth expected", "declining market"),
        ("high customer satisfaction", "many complaints"),
         
    ]
    
    for i, text in enumerate(text_data):
        for claim1, claim2 in claims:
            if claim1 in text and claim2 in text:
                contradictions.append(f"Contradictory claims on slide {i + 1}: '{claim1}' vs '{claim2}'")
    
    return contradictions

def check_timeline_mismatches(text_data):
    """Identify timeline mismatches based on dates found in the text."""
    timeline_mismatches = []
    date_pattern = r'\b\d{1,2}[-/]\d{1,2}[-/]\d{2,4}\b'
    
    dates = {}
    
    for i, text in enumerate(text_data):
        found_dates = re.findall(date_pattern, text)
        for date in found_dates:
            if date in dates:
                dates[date].append(i + 1)
            else:
                dates[date] = [i + 1]
    
    for date, slides in dates.items():
        if len(slides) > 1:
            timeline_mismatches.append(f"Timeline mismatch for date '{date}' found on slides: {slides}")
    
    return timeline_mismatches

def analyze_inconsistencies(text_data):
    """Analyze the text data for various inconsistencies."""
    inconsistencies = []
    inconsistencies.extend(check_conflicting_numerical_data(text_data))
    inconsistencies.extend(check_contradictory_claims(text_data))
    inconsistencies.extend(check_timeline_mismatches(text_data))
    return inconsistencies

def main(pptx_file, image_files):
    """Main function to extract text and analyze inconsistencies."""
    text_data = extract_text_from_pptx(pptx_file)
    for image_file in image_files:
        text_data.append(extract_text_from_image(image_file))
    
    inconsistencies = analyze_inconsistencies(text_data)
    if inconsistencies:
        print("Inconsistencies found:")
        for inconsistency in inconsistencies:
            print(f" - {inconsistency}")
    else:
        print("No inconsistencies found.")

if __name__ == "__main__":
    pptx_file = "D:\\Downloads\\NoogatAssignment\\NoogatAssignment\\NoogatAssignment.pptx"   
    image_files = [
        "D:\\Downloads\\NoogatAssignment\\NoogatAssignment\\NoogatAssignment\\slide1.jpeg", 
        "D:\\Downloads\\NoogatAssignment\\NoogatAssignment\\NoogatAssignment\\slide2.jpeg"
    ]  
    main(pptx_file, image_files)
